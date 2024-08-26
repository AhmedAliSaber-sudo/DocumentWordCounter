import os
import docx2txt
from openpyxl import load_workbook, Workbook
import re
from zipfile import BadZipFile
import win32com.client
from pptx import Presentation
import pythoncom
from datetime import datetime


def count_words_in_file(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    try:
        if file_extension == '.docx':
            text = docx2txt.process(file_path)
        elif file_extension == '.doc':
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch("Word.Application")
            try:
                doc = word.Documents.Open(file_path)
                text = doc.Content.Text
                doc.Close()
            finally:
                word.Quit()
                pythoncom.CoUninitialize()
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            text = " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
        elif file_extension == '.xlsx':
            wb = load_workbook(file_path, read_only=True, data_only=True)
            text = " ".join(
                str(cell.value) for sheet in wb.worksheets for row in sheet.rows for cell in row if cell.value)
            wb.close()
        else:
            print(f"Unsupported file type: {file_path}")
            return None

        words = re.findall(r'\S+', text)
        return len(words)
    except BadZipFile:
        print(f"Error: {file_path} is not a valid Office file. Skipping...")
    except Exception as e:
        print(f"Error processing {file_path}: {str(e)}. Skipping...")

    return None


def list_docs_and_word_counts(folder_path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Document Names and Word Counts"

    ws.append(["Document Name", "Word Count", "Date Modified"])

    file_extensions = ('.docx', '.doc', '.pptx', '.xlsx')
    files = [f for f in os.listdir(folder_path) if f.lower().endswith(file_extensions)]
    files.sort(key=lambda x: os.path.getmtime(os.path.join(folder_path, x)))

    for filename in files:
        file_path = os.path.join(folder_path, filename)
        word_count = count_words_in_file(file_path)
        if word_count is not None:
            date_modified = datetime.fromtimestamp(os.path.getmtime(file_path))
            ws.append([filename, word_count, date_modified])

    wb.save("Document_Names_and_Word_Counts.xlsx")


# Specify the folder path
folder_path = r"G:\documents\work\translation\2024_7"
list_docs_and_word_counts(folder_path)