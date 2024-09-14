import os
import docx2txt
from openpyxl import Workbook
import re
from zipfile import BadZipFile
import win32com.client
from pptx import Presentation
import pythoncom
from datetime import datetime
import tkinter as tk
from tkinter import filedialog, messagebox

class DocumentProcessor:
    @staticmethod
    def count_words_in_docx(file_path):
        text = docx2txt.process(file_path)
        return len(re.findall(r'\S+', text))

    @staticmethod
    def count_words_in_doc(file_path):
        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        try:
            doc = word.Documents.Open(file_path)
            text = doc.Content.Text
            doc.Close()
            return len(re.findall(r'\S+', text))
        finally:
            word.Quit()
            pythoncom.CoUninitialize()

    @staticmethod
    def count_words_in_pptx(file_path):
        prs = Presentation(file_path)
        text = " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
        return len(re.findall(r'\S+', text))

    @staticmethod
    def count_words_in_xlsx(file_path):
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text = " ".join(str(cell.value) for sheet in wb.worksheets for row in sheet.rows for cell in row if cell.value)
        wb.close()
        return len(re.findall(r'\S+', text))

class WordCounter:
    def __init__(self):
        self.processors = {
            '.docx': DocumentProcessor.count_words_in_docx,
            '.doc': DocumentProcessor.count_words_in_doc,
            '.pptx': DocumentProcessor.count_words_in_pptx,
            '.xlsx': DocumentProcessor.count_words_in_xlsx
        }

    def count_words_in_file(self, file_path):
        file_extension = os.path.splitext(file_path)[1].lower()
        processor = self.processors.get(file_extension)
        
        if not processor:
            print(f"Unsupported file type: {file_path}")
            return None

        try:
            return processor(file_path)
        except BadZipFile:
            print(f"Error: {file_path} is not a valid Office file. Skipping...")
        except Exception as e:
            print(f"Error processing {file_path}: {str(e)}. Skipping...")
        
        return None

class ReportGenerator:
    def __init__(self, folder_path, output_path):
        self.folder_path = folder_path
        self.output_path = output_path
        self.word_counter = WordCounter()

    def generate_report(self, status_callback=None):
        wb = Workbook()
        ws = wb.active
        ws.title = "Document Names and Word Counts"
        ws.append(["Document Name", "Word Count", "Date Modified"])

        file_extensions = tuple(self.word_counter.processors.keys())
        files = [f for f in os.listdir(self.folder_path) if f.lower().endswith(file_extensions)]
        files.sort(key=lambda x: os.path.getmtime(os.path.join(self.folder_path, x)))

        total_files = len(files)
        for index, filename in enumerate(files, 1):
            file_path = os.path.join(self.folder_path, filename)
            word_count = self.word_counter.count_words_in_file(file_path)
            if word_count is not None:
                date_modified = datetime.fromtimestamp(os.path.getmtime(file_path))
                ws.append([filename, word_count, date_modified])

            if status_callback:
                status_callback(f"Processing: {index}/{total_files}")

        wb.save(self.output_path)
        return self.output_path

class Application(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Document Word Counter")
        self.geometry("500x300")
        self.folder_path = tk.StringVar()
        self.output_path = tk.StringVar()
        self.create_widgets()

    def create_widgets(self):
        tk.Label(self, text="Select Input Folder:").pack(pady=(10, 0))
        tk.Entry(self, textvariable=self.folder_path, width=50).pack()
        tk.Button(self, text="Browse Input", command=self.browse_input_folder).pack(pady=(0, 10))

        tk.Label(self, text="Select Output File:").pack(pady=(10, 0))
        tk.Entry(self, textvariable=self.output_path, width=50).pack()
        tk.Button(self, text="Browse Output", command=self.browse_output_file).pack(pady=(0, 10))

        self.status_label = tk.Label(self, text="")
        self.status_label.pack(pady=10)

        tk.Button(self, text="Process Files", command=self.process_files).pack(pady=10)

    def browse_input_folder(self):
        folder_selected = filedialog.askdirectory()
        self.folder_path.set(folder_selected)

    def browse_output_file(self):
        file_selected = filedialog.asksaveasfilename(defaultextension=".xlsx",
                                                     filetypes=[("Excel files", "*.xlsx")])
        self.output_path.set(file_selected)

    def process_files(self):
        input_folder = self.folder_path.get()
        output_file = self.output_path.get()

        if not input_folder or not output_file:
            messagebox.showerror("Error", "Please select both input folder and output file")
            return

        self.status_label.config(text="Processing...")
        self.update()

        try:
            report_generator = ReportGenerator(input_folder, output_file)
            output_path = report_generator.generate_report(self.update_status)
            self.status_label.config(text="Processing complete!")
            messagebox.showinfo("Success", f"Excel file created:\n{output_path}")
        except Exception as e:
            self.status_label.config(text="An error occurred")
            messagebox.showerror("Error", str(e))

    def update_status(self, message):
        self.status_label.config(text=message)
        self.update()

if __name__ == "__main__":
    app = Application()
    app.mainloop()